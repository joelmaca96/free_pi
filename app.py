"""GUI ligera (Streamlit) centrada en el Baskonia (primer equipo de config.TEAMS).

Pensada para usuarios sin conocimientos técnicos, sin usar la terminal:
- Resumen: estadísticas avanzadas y forma reciente del Baskonia.
- Partidos anteriores: elegir cualquier partido ya jugado y ver su box score.
- Próximos enfrentamientos: calendario pendiente; al elegir un rival del que
  aun no hay datos, un botón permite descargarlos bajo demanda (roster,
  calendario y box scores recientes) respetando el rate-limit de BBR.
- Plantilla: mosaico de fotos de la plantilla actual (baskonia.com) y ficha
  de cada jugador (posición, dorsal, forma reciente y de temporada).

Solo la descarga bajo demanda de un rival hace peticiones de red; el resto
de la app solo lee `data/baskonia.db`.

Uso:
    streamlit run app.py
"""
import io
import unicodedata
from datetime import datetime, timedelta
from pathlib import Path
from typing import Dict

import pandas as pd
import requests
import streamlit as st
from fpdf import FPDF
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

import config
from db import models
from insights import (
    ZSCORE_COLD_THRESHOLD,
    ZSCORE_HOT_THRESHOLD,
    current_season,
    league_label,
    list_leagues,
    list_seasons,
    parse_minutes,
    per_36,
    player_form_zscore,
    player_load,
    player_recent_form,
    project_next_matchup,
    schedule_difficulty,
    scouting_narrative,
    season_label,
    season_start_year,
    team_advanced_summary,
    validate_data,
)
from main import fetch_opponent_scouting
from scraper.client import BBRClient

st.set_page_config(page_title="Baskonia Pipeline — Informe de partidos", page_icon="🏀", layout="wide")

LOGOS_DIR = Path(__file__).parent / "assets" / "logos"
LOGO_EXTENSIONS = (".png", ".jpg", ".jpeg", ".svg")

# Nº de enfrentamientos directos recientes a mostrar al scoutear un próximo
# rival (independiente de "últimos N partidos" de forma, que es sobre la
# forma general del rival, no específica contra este equipo).
H2H_LAST_N = 2

_WEEKDAYS_ES = {"Mon": "lunes", "Tue": "martes", "Wed": "miércoles", "Thu": "jueves", "Fri": "viernes", "Sat": "sábado", "Sun": "domingo"}
_MONTHS_ES = {
    "Jan": "enero", "Feb": "febrero", "Mar": "marzo", "Apr": "abril", "May": "mayo", "Jun": "junio",
    "Jul": "julio", "Aug": "agosto", "Sep": "septiembre", "Oct": "octubre", "Nov": "noviembre", "Dec": "diciembre",
}


def parse_bbr_date(value: str) -> "datetime | None":
    """Parsea una fecha en formato BBR ('Sun, Nov 23, 2025'), o None si no cuadra."""
    try:
        return datetime.strptime(value, "%a, %b %d, %Y")
    except (ValueError, TypeError):
        return None


def format_date_es(value: str) -> str:
    """Convierte una fecha de BBR a castellano ('domingo, 23 de noviembre de 2025')."""
    dt = parse_bbr_date(value)
    if dt is None:
        return value
    return f"{_WEEKDAYS_ES[dt.strftime('%a')]}, {dt.day} de {_MONTHS_ES[dt.strftime('%b')]} de {dt.year}"


def _fmt(value) -> str:
    return f"{value:.1f}" if value is not None else "-"


def _fmt_pct(value) -> str:
    return f"{value * 100:.1f}%" if value is not None else "-"


def team_logo_path(slug: str) -> "Path | None":
    """Ruta al escudo de un equipo (assets/logos/<slug>.*), o None si no existe."""
    for ext in LOGO_EXTENSIONS:
        path = LOGOS_DIR / f"{slug}{ext}"
        if path.exists():
            return path
    return None


def show_team_logo(slug: str, width: int = 48) -> None:
    """Muestra el escudo de un equipo, o un icono de baloncesto si no hay imagen."""
    path = team_logo_path(slug)
    if path is not None:
        st.image(str(path), width=width)
    else:
        st.markdown(f"<div style='font-size:{width}px; line-height:1'>🏀</div>", unsafe_allow_html=True)


@st.cache_resource
def get_session():
    """Sesión de base de datos compartida entre recargas de la app."""
    Session = models.init_db()
    return Session()


def _team_games(session, team: models.Team, season: "int | None" = None, league: "str | None" = None):
    """Todos los partidos de un equipo (jugados y pendientes), ordenados por fecha.

    Punto único de acceso al calendario de un equipo: aquí se aplican los dos
    filtros globales de la app. `league` es una columna real (`Game.league`), así
    que se filtra en la propia consulta; `season` es derivada de la fecha (ver
    `insights.season_start_year`), así que se filtra en Python tras traer las
    filas. `None` en cualquiera de los dos = sin filtrar por ese eje.
    """
    query = session.query(models.Game).filter(
        (models.Game.home_team_id == team.id) | (models.Game.away_team_id == team.id)
    )
    if league is not None:
        query = query.filter(models.Game.league == league)
    games = query.all()
    if season is not None:
        games = [g for g in games if season_start_year(g.date) == season]
    games.sort(key=lambda g: parse_bbr_date(g.date) or datetime.min)
    return games


def _team_stats_for_game(session, game_id: int, team_id: int):
    return session.query(models.TeamGameStats).filter_by(game_id=game_id, team_id=team_id).first()


def _rival_of(game: models.Game, team: models.Team) -> models.Team:
    """Devuelve el equipo rival de `team` en un partido dado."""
    return game.away_team if game.home_team_id == team.id else game.home_team


def _result_label(game: models.Game, team: models.Team) -> str:
    """Resultado del partido desde el punto de vista de `team`, o 'pendiente' si no se ha jugado."""
    is_home = game.home_team_id == team.id
    team_score = game.home_score if is_home else game.away_score
    opp_score = game.away_score if is_home else game.home_score
    if team_score is None:
        return "pendiente"
    return f"{team_score}-{opp_score}"


def past_games(session, team: models.Team, season: "int | None" = None, league: "str | None" = None) -> list:
    """Partidos ya jugados de un equipo (con resultado), del más reciente al más antiguo.

    `season`/`league` acotan a la temporada/competición seleccionadas; `None` en
    cualquiera de los dos no filtra por ese eje.
    """
    query = (
        session.query(models.Game)
        .filter((models.Game.home_team_id == team.id) | (models.Game.away_team_id == team.id))
        .filter(models.Game.home_score.isnot(None))
    )
    if league is not None:
        query = query.filter(models.Game.league == league)
    games = query.all()
    if season is not None:
        games = [g for g in games if season_start_year(g.date) == season]
    games.sort(key=lambda g: parse_bbr_date(g.date) or datetime.min, reverse=True)
    return games


def upcoming_games(session, team: models.Team) -> list:
    """Próximos partidos de un equipo (sin resultado todavía), del más cercano al más lejano.

    Un partido sin resultado no siempre es "próximo": BBR deja para siempre
    sin resultado la fila de un partido aplazado, aunque se haya jugado más
    adelante en otra fecha como fila aparte. Se excluyen las fechas ya
    pasadas (por si hay algún caso similar sin anotar), para no mezclar
    partidos de la temporada ya cerrada con los realmente pendientes.
    """
    games = (
        session.query(models.Game)
        .filter((models.Game.home_team_id == team.id) | (models.Game.away_team_id == team.id))
        .filter(models.Game.home_score.is_(None))
        .all()
    )
    today = datetime.now()
    games = [g for g in games if (parse_bbr_date(g.date) or datetime.max) >= today]
    games.sort(key=lambda g: parse_bbr_date(g.date) or datetime.max)
    return games


def _games_to_df(session, games: list, team: models.Team) -> pd.DataFrame:
    """Tabla de una lista concreta de partidos de un equipo, con pace/ORtg/DRtg/Net."""
    rows = []
    for game in games:
        is_home = game.home_team_id == team.id
        rival = game.away_team.name if is_home else game.home_team.name
        team_score = game.home_score if is_home else game.away_score
        opp_score = game.away_score if is_home else game.home_score
        stats = _team_stats_for_game(session, game.id, team.id)
        rows.append(
            {
                "Fecha": format_date_es(game.date),
                "Rival": rival,
                "Resultado": f"{team_score}-{opp_score}" if team_score is not None else "-",
                "Pace": round(stats.pace, 1) if stats and stats.pace is not None else None,
                "ORtg": round(stats.off_rating, 1) if stats and stats.off_rating is not None else None,
                "DRtg": round(stats.def_rating, 1) if stats and stats.def_rating is not None else None,
                "Net": round(stats.net_rating, 1) if stats and stats.net_rating is not None else None,
            }
        )
    return pd.DataFrame(rows)


def team_summary_df(
    session, team: models.Team, season: "int | None" = None, league: "str | None" = None
) -> pd.DataFrame:
    """Tabla de todos los partidos guardados de un equipo (jugados y pendientes).

    El filtro es por bucket temporal (fecha del partido) y competición, no por si
    ya se jugó: un partido pendiente cuya fecha cae dentro de la temporada
    seleccionada sigue apareciendo, con el resultado a "-".
    """
    return _games_to_df(session, _team_games(session, team, season, league), team)


def recent_games_df(
    session, team: models.Team, last_n: int, season: "int | None" = None, league: "str | None" = None
) -> pd.DataFrame:
    """Tabla de los últimos `last_n` partidos JUGADOS de un equipo, sea quien sea el rival."""
    played = [g for g in _team_games(session, team, season, league) if g.home_score is not None]
    return _games_to_df(session, played[-last_n:], team)


def head_to_head_summary_df(
    session, team: models.Team, season: "int | None" = None, league: "str | None" = None
) -> pd.DataFrame:
    """Tabla de los partidos jugados contra los otros equipos de `config.TEAMS`.

    Separada de `recent_games_df` a propósito: un enfrentamiento directo puede
    haberse jugado hace tiempo y quedar fuera de los últimos N partidos, así
    que mezclarlos en una sola tabla puede dar la impresión equivocada de que
    no hay enfrentamientos directos recientes cuando en realidad no se han
    mirado los partidos anteriores a los últimos N.
    """
    rival_slugs = {slug for slug in config.TEAMS if slug != team.slug}
    played = [g for g in _team_games(session, team, season, league) if g.home_score is not None]
    games = [g for g in played if _rival_of(g, team).slug in rival_slugs]
    return _games_to_df(session, games, team)


def recent_form_df(
    session, team: models.Team, last_n: int, season: "int | None" = None, league: "str | None" = None
) -> pd.DataFrame:
    """Tabla de forma reciente (medias últimos N partidos) por jugador.

    "3PA%" es el porcentaje de intentos de tiro que son triples y "FTr" la tasa
    de tiros libres (FTA/FGA, no el % de acierto): perfil de tiro del jugador.
    """
    form = player_recent_form(session, team, last_n=last_n, season=season, league=league)
    rows = [
        {
            "Jugador": r["player_name"],
            "PJ": r["games"],
            "MIN": round(r["avg_minutes"], 1),
            "PTS": round(r["avg_pts"], 1),
            "PTS/36": round(r["avg_pts_per36"], 1) if r["avg_pts_per36"] is not None else None,
            "eFG%": round(r["avg_efg_pct"] * 100, 1) if r["avg_efg_pct"] is not None else None,
            "TS%": round(r["avg_ts_pct"] * 100, 1) if r["avg_ts_pct"] is not None else None,
            "3PA%": round(r["fg3a_rate"] * 100, 1) if r["fg3a_rate"] is not None else None,
            "FTr": round(r["ft_rate"], 2) if r["ft_rate"] is not None else None,
        }
        for r in form
    ]
    return pd.DataFrame(rows)


def _streak_label(z_score: "float | None") -> str:
    """Etiqueta de racha para un z-score (umbrales de `insights`), o '-' si no hay dato."""
    if z_score is None:
        return "-"
    if z_score >= ZSCORE_HOT_THRESHOLD:
        return "🔥 en racha"
    if z_score <= ZSCORE_COLD_THRESHOLD:
        return "❄️ bajo forma"
    return "➖"


def streaks_df(
    session, team: models.Team, season: int, recent_n: int, league: "str | None" = None
) -> pd.DataFrame:
    """Tabla de rachas por jugador dentro de una temporada: doble z-score (PTS y TS%).

    Las dos columnas de racha se mantienen separadas (volumen y eficiencia) en vez
    de fusionarse en una sola etiqueta: son señales independientes.
    """
    rows = []
    for r in player_form_zscore(session, team, season=season, recent_n=recent_n, league=league):
        rows.append(
            {
                "Jugador": r["player_name"],
                "PJ temporada": r["games_season"],
                f"PTS últimos {recent_n}": round(r["recent_avg_pts"], 1) if r["recent_avg_pts"] is not None else None,
                "PTS temporada": round(r["season_avg_pts"], 1) if r["season_avg_pts"] is not None else None,
                "z-score PTS": round(r["z_score_pts"], 2) if r["z_score_pts"] is not None else None,
                "Racha PTS": _streak_label(r["z_score_pts"]),
                f"TS% últimos {recent_n}": round(r["recent_avg_ts_pct"] * 100, 1) if r["recent_avg_ts_pct"] is not None else None,
                "TS% temporada": round(r["season_avg_ts_pct"] * 100, 1) if r["season_avg_ts_pct"] is not None else None,
                "z-score TS%": round(r["z_score_ts"], 2) if r["z_score_ts"] is not None else None,
                "Racha TS%": _streak_label(r["z_score_ts"]),
            }
        )
    return pd.DataFrame(rows)


def games_in_window(
    session, team: models.Team, window_days: int, reference_date: "datetime | None" = None
) -> list:
    """Partidos ya JUGADOS de `team` dentro de los últimos `window_days`
    días respecto a `reference_date` (por defecto, ahora).

    Se llama a `past_games` sin temporada ni competición a propósito: la ventana
    de días ya es un filtro temporal más preciso, y la carga física de un jugador
    es transversal a la competición (un partido de ACB y otro de Euroliga en la
    misma semana cansan igual).
    """
    reference_date = reference_date or datetime.now()
    cutoff = reference_date - timedelta(days=window_days)
    return [
        g for g in past_games(session, team)
        if cutoff <= (parse_bbr_date(g.date) or datetime.min) <= reference_date
    ]


def player_load_df(session, team: models.Team, window_days: int) -> pd.DataFrame:
    """Tabla de carga de minutos por jugador en los últimos `window_days` días."""
    games = games_in_window(session, team, window_days)
    rows = [
        {
            "Jugador": r["player_name"],
            "PJ ventana": r["games"],
            "MIN totales": round(r["total_minutes"], 1),
            "MIN/partido": round(r["avg_minutes"], 1),
        }
        for r in player_load(session, team, games)
    ]
    return pd.DataFrame(rows)


def schedule_difficulty_df(difficulty: Dict[str, object]) -> pd.DataFrame:
    """Tabla de los próximos rivales de una salida de `insights.schedule_difficulty`."""
    rows = [
        {
            "Fecha": format_date_es(o["date"]),
            "Rival": o["opponent_name"],
            "Net Rating": round(o["net_rating"], 1) if o["net_rating"] is not None else None,
        }
        for o in difficulty["opponents"]
    ]
    return pd.DataFrame(rows)


def boxscore_df(session, game: models.Game, team: models.Team) -> pd.DataFrame:
    """Tabla de box score (PTS, REB, AST, PTS/36, eFG%, TS%) de un equipo en un partido."""
    rows = (
        session.query(models.BoxScore)
        .filter_by(game_id=game.id, team_id=team.id)
        .order_by(models.BoxScore.points.desc())
        .all()
    )
    data = []
    for row in rows:
        minutes = parse_minutes(row.minutes)
        pts_per36 = per_36(row.points, minutes)
        data.append(
            {
                "Jugador": row.player_name,
                "MIN": row.minutes or pd.NA,
                "PTS": row.points if row.points is not None else pd.NA,
                "REB": row.rebounds if row.rebounds is not None else pd.NA,
                "AST": row.assists if row.assists is not None else pd.NA,
                "PTS/36": round(pts_per36, 1) if pts_per36 is not None else pd.NA,
                "eFG%": round(row.efg_pct * 100, 1) if row.efg_pct is not None else pd.NA,
                "TS%": round(row.ts_pct * 100, 1) if row.ts_pct is not None else pd.NA,
            }
        )
    return pd.DataFrame(data)


def head_to_head_games(
    session,
    team_a: models.Team,
    team_b: models.Team,
    season: "int | None" = None,
    league: "str | None" = None,
):
    """Enfrentamientos directos entre dos equipos, ordenados por fecha.

    No distingue jugado/pendiente (nunca lo hizo): un enfrentamiento ya
    programado de la temporada seleccionada aparece con el marcador sin rellenar.
    """
    query = session.query(models.Game).filter(
        ((models.Game.home_team_id == team_a.id) & (models.Game.away_team_id == team_b.id))
        | ((models.Game.home_team_id == team_b.id) & (models.Game.away_team_id == team_a.id))
    )
    if league is not None:
        query = query.filter(models.Game.league == league)
    games = query.all()
    if season is not None:
        games = [g for g in games if season_start_year(g.date) == season]
    games.sort(key=lambda g: parse_bbr_date(g.date) or datetime.min)
    return games


def render_narrative_section(
    session, team: models.Team, season: int, recent_n: int, league: "str | None" = None
) -> None:
    """Subsección 'Resumen automático': narrativa de estilo y forma del equipo.

    No pinta nada si el equipo no tiene estadísticas avanzadas en la
    temporada/competición seleccionadas (nada que resumir).
    """
    narrative = scouting_narrative(session, team, season=season, recent_n=recent_n, league=league)
    if narrative is None:
        return
    st.subheader("Resumen automático")
    st.write(narrative)


def render_streaks_section(
    session, team: models.Team, season: int, recent_n: int, league: "str | None" = None
) -> None:
    """Subsección 'Rachas (hot/cold)': tabla con doble z-score (PTS y TS%) por jugador."""
    st.subheader("Rachas (hot/cold)")
    st.caption(
        f"Z-score de los últimos {recent_n} partidos de cada jugador frente a su propia media "
        f"de la temporada {season_label(season)}: volumen (PTS) y eficiencia (TS%) por separado."
    )
    df_streaks = streaks_df(session, team, season, recent_n, league)
    if df_streaks.empty:
        st.info(
            f"Sin jugadores con partidos suficientes en la temporada {season_label(season)} "
            "para calcular racha todavía."
        )
    else:
        st.dataframe(df_streaks.fillna(pd.NA), use_container_width=True, hide_index=True)


def render_player_load_section(session, team: models.Team) -> None:
    """Subsección 'Carga de minutos': minutos acumulados por jugador en una ventana de días."""
    st.subheader("Carga de minutos (gestión de fatiga)")
    window_days = st.number_input(
        "Ventana de días",
        min_value=1,
        max_value=30,
        value=14,
        key=f"load_window_{team.id}",
    )
    st.caption(
        "Minutos acumulados en los últimos días naturales, independientemente de la temporada y "
        "la competición seleccionadas: la carga física de un jugador es transversal a ambas."
    )
    df_load = player_load_df(session, team, int(window_days))
    if df_load.empty:
        st.info(f"Sin partidos jugados con minutos registrados en los últimos {int(window_days)} días.")
    else:
        st.dataframe(df_load.fillna(pd.NA), use_container_width=True, hide_index=True)


def render_team_tab(
    session, team: models.Team, last_n: int, season: int, league: "str | None" = None
) -> None:
    """Contenido de la pestaña de un equipo: estadísticas avanzadas, partidos y forma reciente."""
    logo_col, title_col, pdf_col = st.columns([1, 7, 2])
    with logo_col:
        show_team_logo(team.slug, width=56)
    with title_col:
        st.subheader(team.name)
    with pdf_col:
        rival_slug = next((s for s in config.TEAMS if s != team.slug), None)
        rival = session.query(models.Team).filter_by(slug=rival_slug).first() if rival_slug else None
        if rival is not None:
            pdf_bytes = build_pdf_report(session, team, rival, last_n, season, league)
            st.download_button(
                "📄 Generar informe en PDF",
                data=pdf_bytes,
                file_name=f"informe_{team.slug}_vs_{rival.slug}.pdf",
                mime="application/pdf",
                key=f"pdf_resumen_{team.id}",
            )

    render_narrative_section(session, team, season, last_n, league)

    st.subheader("Estadísticas avanzadas (medias)")
    summary = team_advanced_summary(session, team, season=season, league=league)
    m1, m2, m3, m4, m5, m6 = st.columns(6)
    m1.metric("Pace", _fmt(summary["avg_pace"]))
    m2.metric("ORtg", _fmt(summary["avg_off_rating"]))
    m3.metric("DRtg", _fmt(summary["avg_def_rating"]))
    m4.metric("Net Rating", _fmt(summary["avg_net_rating"]))
    m5.metric("eFG%", _fmt_pct(summary["avg_efg_pct"]))
    m6.metric("TS%", _fmt_pct(summary["avg_ts_pct"]))

    st.subheader(f"Últimos {last_n} partidos jugados")
    df_recent = recent_games_df(session, team, last_n, season, league)
    if df_recent.empty:
        st.info("Sin partidos guardados todavía.")
    else:
        chart_df = df_recent.dropna(subset=["ORtg", "DRtg"])
        if not chart_df.empty:
            st.bar_chart(chart_df.set_index("Fecha")[["ORtg", "DRtg"]])
        st.dataframe(df_recent.fillna(pd.NA), use_container_width=True, hide_index=True)

    other_teams = [slug for slug in config.TEAMS if slug != team.slug]
    if other_teams:
        st.subheader("Enfrentamientos directos")
        st.caption(
            "Partidos contra los otros equipos de interés configurados (TEAMS). "
            "Pueden solaparse con los últimos partidos jugados si el enfrentamiento "
            "es reciente, o quedar fuera de ellos si fue hace más tiempo."
        )
        df_h2h = head_to_head_summary_df(session, team, season, league)
        if df_h2h.empty:
            st.info("Sin enfrentamientos directos jugados todavía.")
        else:
            st.dataframe(df_h2h.fillna(pd.NA), use_container_width=True, hide_index=True)

    st.subheader(f"Forma reciente (últimos {last_n} partidos jugados)")
    df_form = recent_form_df(session, team, last_n, season, league)
    if df_form.empty:
        st.info("Sin datos suficientes.")
    else:
        st.bar_chart(df_form.set_index("Jugador")["PTS"])
        st.dataframe(df_form.fillna(pd.NA), use_container_width=True, hide_index=True)

    render_streaks_section(session, team, season, last_n, league)
    render_player_load_section(session, team)


def render_head_to_head_tab(
    session,
    team_a: models.Team,
    team_b: models.Team,
    last_n: "int | None" = None,
    season: "int | None" = None,
    league: "str | None" = None,
) -> None:
    """Contenido de la pestaña de enfrentamientos directos.

    `last_n` limita a los últimos N enfrentamientos jugados (los más
    recientes) en vez de todo el histórico; `None` los muestra todos.
    """
    games = head_to_head_games(session, team_a, team_b, season, league)
    if not games:
        st.info("Sin enfrentamientos directos guardados todavía entre estos dos equipos.")
        return
    if last_n is not None:
        games = games[-last_n:]

    for game in games:
        logo_col, title_col = st.columns([1, 8])
        with logo_col:
            show_team_logo(game.home_team.slug, width=32)
        with title_col:
            st.markdown(
                f"### {format_date_es(game.date)} — {game.home_team.name} {game.home_score} - "
                f"{game.away_score} {game.away_team.name}"
            )
        stats_home = _team_stats_for_game(session, game.id, game.home_team_id)
        stats_away = _team_stats_for_game(session, game.id, game.away_team_id)
        if stats_home and stats_away:
            m1, m2, m3 = st.columns(3)
            m1.metric("Pace", round(stats_home.pace, 1) if stats_home.pace is not None else "-")
            m2.metric(f"Net Rating {game.home_team.name}", round(stats_home.net_rating, 1) if stats_home.net_rating is not None else "-")
            m3.metric(f"Net Rating {game.away_team.name}", round(stats_away.net_rating, 1) if stats_away.net_rating is not None else "-")

        col_home, col_away = st.columns(2)
        with col_home:
            show_team_logo(game.home_team.slug, width=28)
            st.markdown(f"**{game.home_team.name}**")
            st.dataframe(boxscore_df(session, game, game.home_team), use_container_width=True, hide_index=True)
        with col_away:
            show_team_logo(game.away_team.slug, width=28)
            st.markdown(f"**{game.away_team.name}**")
            st.dataframe(boxscore_df(session, game, game.away_team), use_container_width=True, hide_index=True)
        st.divider()


def render_past_games_tab(
    session, team: models.Team, last_n: int, season: int, league: "str | None" = None
) -> None:
    """Permite elegir cualquier partido anterior de Baskonia y ver su box score."""
    games = past_games(session, team, season, league)
    if not games:
        st.info("Sin partidos guardados todavía. Ejecuta `python main.py`.")
        return

    labels = {g.id: f"{format_date_es(g.date)} — {_rival_of(g, team).name} ({_result_label(g, team)})" for g in games}
    game_id = st.selectbox("Partido", [g.id for g in games], format_func=lambda i: labels[i])
    game = next(g for g in games if g.id == game_id)
    rival = _rival_of(game, team)

    logo_col, title_col = st.columns([1, 8])
    with logo_col:
        show_team_logo(rival.slug, width=48)
    with title_col:
        st.markdown(
            f"### {format_date_es(game.date)} — {game.home_team.name} {game.home_score} - "
            f"{game.away_score} {game.away_team.name}"
        )

    stats_home = _team_stats_for_game(session, game.id, game.home_team_id)
    stats_away = _team_stats_for_game(session, game.id, game.away_team_id)
    if stats_home and stats_away:
        m1, m2, m3 = st.columns(3)
        m1.metric("Pace", _fmt(stats_home.pace))
        m2.metric(f"Net Rating {game.home_team.name}", _fmt(stats_home.net_rating))
        m3.metric(f"Net Rating {game.away_team.name}", _fmt(stats_away.net_rating))

    col_home, col_away = st.columns(2)
    with col_home:
        show_team_logo(game.home_team.slug, width=28)
        st.markdown(f"**{game.home_team.name}**")
        st.dataframe(boxscore_df(session, game, game.home_team), use_container_width=True, hide_index=True)
    with col_away:
        show_team_logo(game.away_team.slug, width=28)
        st.markdown(f"**{game.away_team.name}**")
        st.dataframe(boxscore_df(session, game, game.away_team), use_container_width=True, hide_index=True)

    pdf_bytes = build_pdf_report(session, team, rival, last_n, season, league)
    st.download_button(
        "📄 Informe de este partido en PDF",
        data=pdf_bytes,
        file_name=f"informe_{team.slug}_vs_{rival.slug}.pdf",
        mime="application/pdf",
        key=f"pdf_past_{game.id}",
    )


def render_schedule_difficulty_section(
    session, team: models.Team, season: int, upcoming: list, league: "str | None" = None
) -> None:
    """Subsección 'Dificultad del próximo tramo': Net Rating medio de los próximos N rivales."""
    st.subheader("Dificultad del próximo tramo de calendario")
    next_n = st.number_input(
        "Próximos N partidos",
        min_value=1,
        max_value=15,
        value=5,
        key=f"next_n_{team.id}",
    )
    difficulty = schedule_difficulty(
        session, team, upcoming, season=season, next_n=int(next_n), league=league
    )
    if difficulty["games_considered"] == 0:
        st.info(f"Sin partidos pendientes de {league_label(league)} en el calendario descargado.")
        return

    m1, m2, m3 = st.columns(3)
    m1.metric("Partidos considerados", difficulty["games_considered"])
    m2.metric("Rivales con datos", f"{difficulty['opponents_scouted']}/{difficulty['games_considered']}")
    m3.metric("Net Rating medio del rival", _fmt(difficulty["avg_opponent_net_rating"]))
    st.caption(
        f"Net Rating de cada rival según su temporada {season_label(season)} "
        f"({league_label(league)}); cuanto más alto, más difícil el tramo."
    )
    st.dataframe(
        schedule_difficulty_df(difficulty).fillna(pd.NA), use_container_width=True, hide_index=True
    )
    if difficulty["avg_opponent_net_rating"] is None:
        st.info(
            "Ningún rival de este tramo tiene estadísticas avanzadas guardadas en esa "
            "temporada/competición todavía."
        )


def render_matchup_projection_section(
    session, team: models.Team, rival: models.Team, season: int, league: "str | None" = None
) -> None:
    """Subsección 'Proyección del partido': marcador esperado a partir de pace y ratings."""
    st.subheader("Proyección del partido")
    projection = project_next_matchup(session, team, rival, season=season, league=league)
    if projection is None:
        st.info(
            "Datos insuficientes para proyectar el marcador: falta pace/ORtg/DRtg de alguno de "
            "los dos equipos en la temporada y competición seleccionadas."
        )
        return

    m1, m2, m3 = st.columns(3)
    m1.metric("Posesiones proyectadas", _fmt(projection["projected_possessions"]))
    m2.metric(f"{team.name} (proyección)", _fmt(projection["team_projected_score"]))
    m3.metric(f"{rival.name} (proyección)", _fmt(projection["opp_projected_score"]))
    st.caption(
        "Estimación simple: posesiones = media de los dos paces; rating esperado de cada equipo = "
        "media entre su ataque y la defensa del rival."
    )


def render_upcoming_tab(
    session, team: models.Team, last_n: int, season: int, league: "str | None" = None
) -> None:
    """Lista el calendario pendiente y permite scoutear bajo demanda al rival elegido."""
    games = upcoming_games(session, team)
    if not games:
        st.info(
            "No hay partidos pendientes en el calendario descargado. Puede que Basketball-Reference todavía no "
            "haya publicado el calendario de la próxima temporada."
        )
        return

    render_schedule_difficulty_section(session, team, season, games, league)

    labels = {g.id: f"{format_date_es(g.date)} — {_rival_of(g, team).name}" for g in games}
    game_id = st.selectbox("Próximo enfrentamiento", [g.id for g in games], format_func=lambda i: labels[i])
    game = next(g for g in games if g.id == game_id)
    rival = _rival_of(game, team)

    logo_col, title_col = st.columns([1, 8])
    with logo_col:
        show_team_logo(rival.slug, width=48)
    with title_col:
        lugar = "en casa" if game.home_team_id == team.id else "fuera"
        st.markdown(f"### {format_date_es(game.date)} — {rival.name} ({lugar})")

    has_roster = session.query(models.Player).filter_by(team_id=rival.id).first() is not None
    if not has_roster:
        st.warning(f"Todavía no hay datos de {rival.name} en la base de datos.")
        n_rival = st.number_input(
            f"Últimos N partidos de {rival.name} a descargar",
            min_value=1,
            max_value=15,
            value=5,
            key=f"n_rival_{rival.id}",
        )
        st.caption(
            f"Se respeta el rate-limit de Basketball-Reference (~20s por petición): "
            f"puede tardar del orden de {int((n_rival + 2) * 20)}s."
        )
        if st.button(f"📥 Descargar datos de {rival.name}", key=f"fetch_{rival.id}"):
            with st.spinner(f"Descargando roster, calendario y box scores de {rival.name}..."):
                client = BBRClient()
                try:
                    fetch_opponent_scouting(session, client, rival, int(n_rival))
                except RuntimeError as exc:
                    st.error(str(exc))
                    return
            st.rerun()
        return

    render_matchup_projection_section(session, team, rival, season, league)

    st.subheader(f"Scouting: {rival.name}")
    render_team_tab(session, rival, last_n, season, league)

    st.subheader(f"Últimos {H2H_LAST_N} enfrentamientos directos: {team.name} vs {rival.name}")
    render_head_to_head_tab(session, team, rival, last_n=H2H_LAST_N, season=season, league=league)

    pdf_bytes = build_pdf_report(session, team, rival, last_n, season, league)
    st.download_button(
        "📄 Informe de scouting en PDF",
        data=pdf_bytes,
        file_name=f"scouting_{team.slug}_vs_{rival.slug}.pdf",
        mime="application/pdf",
        key=f"pdf_upcoming_{game.id}",
    )


def current_roster(session, team: models.Team) -> list:
    """Jugadores de la plantilla actual, ordenados por dorsal.

    Solo la plantilla oficial de baskonia.com (ver
    `scraper/baskonia_official.py`) rellena `photo_url`; un jugador
    capturado únicamente vía BBR (temporadas o rosters anteriores, ya no en
    el equipo) no tiene foto y no se considera "plantilla actual".
    """
    players = (
        session.query(models.Player)
        .filter_by(team_id=team.id)
        .filter(models.Player.photo_url.isnot(None))
        .all()
    )
    players.sort(key=lambda p: int(p.number) if p.number and p.number.isdigit() else 999)
    return players


def _player_stats_row(
    session,
    team: models.Team,
    player: models.Player,
    last_n: int,
    season: "int | None" = None,
    league: "str | None" = None,
) -> "dict | None":
    """Fila de `player_recent_form` para un jugador concreto, o `None` si no tiene partidos."""
    form = player_recent_form(session, team, last_n=last_n, season=season, league=league)
    return next((r for r in form if r["player_name"] == player.name), None)


def render_player_card(
    session,
    team: models.Team,
    player: models.Player,
    last_n: int,
    season: "int | None" = None,
    league: "str | None" = None,
) -> None:
    """Ficha de un jugador: foto, posición, dorsal, forma reciente y de temporada."""
    col_photo, col_info = st.columns([1, 3])
    with col_photo:
        st.image(player.photo_url, width=200)
    with col_info:
        st.subheader(player.name)
        st.write(f"**Posición:** {player.position or '-'}")
        st.write(f"**Dorsal:** {player.number or '-'}")

    st.subheader(f"Forma reciente (últimos {last_n} partidos jugados)")
    recent = _player_stats_row(session, team, player, last_n, season=season, league=league)
    if recent is None:
        st.info("Sin datos de partidos recientes todavía (jugador nuevo o sin box scores capturados).")
    else:
        m1, m2, m3, m4, m5, m6, m7 = st.columns(7)
        m1.metric("Partidos", recent["games"])
        m2.metric("MIN", _fmt(recent["avg_minutes"]))
        m3.metric("PTS", _fmt(recent["avg_pts"]))
        m4.metric("eFG%", _fmt_pct(recent["avg_efg_pct"]))
        m5.metric("TS%", _fmt_pct(recent["avg_ts_pct"]))
        m6.metric("3PA%", _fmt_pct(recent["fg3a_rate"]))
        m7.metric("FTr", _fmt(recent["ft_rate"]))

    st.subheader("Estadísticas de la temporada")
    # `last_n=1000` es un límite práctico "sin techo": con `season` explícito, la
    # fila ya queda acotada a la temporada seleccionada.
    season_stats = _player_stats_row(session, team, player, last_n=1000, season=season, league=league)
    if season_stats is None:
        st.info("Sin partidos registrados esta temporada todavía.")
    else:
        m1, m2, m3, m4, m5, m6, m7 = st.columns(7)
        m1.metric("Partidos", season_stats["games"])
        m2.metric("MIN", _fmt(season_stats["avg_minutes"]))
        m3.metric("PTS", _fmt(season_stats["avg_pts"]))
        m4.metric("eFG%", _fmt_pct(season_stats["avg_efg_pct"]))
        m5.metric("TS%", _fmt_pct(season_stats["avg_ts_pct"]))
        m6.metric("3PA%", _fmt_pct(season_stats["fg3a_rate"]))
        m7.metric("FTr", _fmt(season_stats["ft_rate"]))


# Las 4 estadísticas del PPT: clave en player_recent_form, etiqueta, formato
# y si "más alto es mejor" (para decidir el color verde/rojo al comparar
# contra la media del equipo). Pérdidas es la única al revés: menos es mejor.
# (Se descarta +/- pese a estar calculado en insights.py: los box scores
# internacionales de BBR no traen esa columna, así que siempre saldría
# "sin datos" — verificado contra una página de box score real.)
_PPT_STATS = [
    ("avg_pts", "PTS", lambda v: f"{v:.1f}", True),
    ("avg_efg_pct", "eFG%", lambda v: f"{v * 100:.1f}%", True),
    ("avg_ts_pct", "TS%", lambda v: f"{v * 100:.1f}%", True),
    ("avg_turnovers", "Pérdidas", lambda v: f"{v:.1f}", False),
]
_PPT_GOOD_COLOR = RGBColor(0x1B, 0x8A, 0x3A)
_PPT_BAD_COLOR = RGBColor(0xC0, 0x1C, 0x28)
_PPT_NEUTRAL_COLOR = RGBColor(0x66, 0x66, 0x66)


def _fetch_image_bytes(url: str) -> "bytes | None":
    """Descarga una foto y la normaliza a PNG para incrustarla en el PPT.

    python-pptx no soporta WEBP (el formato del icono genérico de fallback de
    baskonia.com para jugadores sin foto subida todavía); se reconvierte con
    Pillow independientemente del formato original en vez de asumir uno.
    Devuelve `None` si falla (no bloquea la generación de esa diapositiva).
    """
    try:
        response = requests.get(url, headers={"User-Agent": config.USER_AGENT}, timeout=10)
        response.raise_for_status()
        image = Image.open(io.BytesIO(response.content)).convert("RGB")
        png_buffer = io.BytesIO()
        image.save(png_buffer, format="PNG")
        return png_buffer.getvalue()
    except Exception:  # noqa: BLE001
        return None


def build_roster_pptx(
    session,
    team: models.Team,
    players: list,
    last_n: int,
    season: "int | None" = None,
    league: "str | None" = None,
) -> bytes:
    """Genera un PPTX con una diapositiva por jugador: foto, nombre y sus 4
    estadísticas más relevantes (PTS, eFG%, TS%, pérdidas; medias de los
    últimos `last_n` partidos jugados).

    Cada estadística se colorea en verde si el jugador está "para bien"
    respecto a la media del equipo en esa estadística (por encima si más es
    mejor, por debajo si es al revés, como en pérdidas) o en rojo si está
    "para mal"; gris si no hay datos suficientes para comparar.
    """
    rows = player_recent_form(session, team, last_n=last_n, season=season, league=league)
    rows_by_name = {r["player_name"]: r for r in rows}

    team_avgs: Dict[str, "float | None"] = {}
    for key, _, _, _ in _PPT_STATS:
        values = [r[key] for r in rows if r.get(key) is not None]
        team_avgs[key] = sum(values) / len(values) if values else None

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for player in players:
        slide = prs.slides.add_slide(blank_layout)

        photo_bytes = _fetch_image_bytes(player.photo_url) if player.photo_url else None
        if photo_bytes:
            slide.shapes.add_picture(io.BytesIO(photo_bytes), Inches(0.5), Inches(0.7), height=Inches(5.5))

        name_box = slide.shapes.add_textbox(Inches(4.3), Inches(0.6), Inches(5.2), Inches(0.9))
        name_run = name_box.text_frame.paragraphs[0].add_run()
        name_run.text = f"#{player.number or '-'} {player.name}"
        name_run.font.size = Pt(30)
        name_run.font.bold = True

        pos_box = slide.shapes.add_textbox(Inches(4.3), Inches(1.35), Inches(5.2), Inches(0.5))
        pos_run = pos_box.text_frame.paragraphs[0].add_run()
        pos_run.text = player.position or "Posición desconocida"
        pos_run.font.size = Pt(16)
        pos_run.font.color.rgb = _PPT_NEUTRAL_COLOR

        row = rows_by_name.get(player.name)
        for i, (key, label, fmt, higher_is_better) in enumerate(_PPT_STATS):
            value = row.get(key) if row else None
            box = slide.shapes.add_textbox(Inches(4.3), Inches(2.2 + i * 0.95), Inches(5.2), Inches(0.85))
            paragraph = box.text_frame.paragraphs[0]

            label_run = paragraph.add_run()
            label_run.text = f"{label}: "
            label_run.font.size = Pt(24)
            label_run.font.bold = True

            value_run = paragraph.add_run()
            value_run.font.size = Pt(24)
            value_run.font.bold = True
            team_avg = team_avgs.get(key)
            if value is None:
                value_run.text = "sin datos"
                value_run.font.color.rgb = _PPT_NEUTRAL_COLOR
            else:
                value_run.text = fmt(value)
                if team_avg is None:
                    value_run.font.color.rgb = _PPT_NEUTRAL_COLOR
                else:
                    is_good = (value >= team_avg) if higher_is_better else (value <= team_avg)
                    value_run.font.color.rgb = _PPT_GOOD_COLOR if is_good else _PPT_BAD_COLOR

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def render_roster_tab(
    session, team: models.Team, last_n: int, season: int, league: "str | None" = None
) -> None:
    """Plantilla actual del equipo: mosaico de fotos y ficha del jugador elegido."""
    players = current_roster(session, team)
    if not players:
        st.info(
            "Sin plantilla descargada todavía. Ejecuta `python main.py` para obtenerla "
            "de la web oficial del Baskonia (baskonia.com)."
        )
        return

    header_col, button_col = st.columns([4, 2])
    with header_col:
        st.subheader(f"Plantilla actual ({len(players)} jugadores)")
    with button_col:
        if st.button("🎬 Generar ppt para Paolo", key="ppt_paolo"):
            with st.spinner("Generando presentación..."):
                pptx_bytes = build_roster_pptx(session, team, players, last_n, season, league)
            st.download_button(
                "⬇️ Descargar presentación (.pptx)",
                data=pptx_bytes,
                file_name=f"plantilla_{team.slug}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                key="download_ppt_paolo",
            )

    cols = st.columns(6)
    for i, p in enumerate(players):
        with cols[i % 6]:
            st.image(p.photo_url, width=100)
            st.caption(f"#{p.number or '-'} {p.name}")

    st.divider()
    labels = {p.id: f"#{p.number or '-'} {p.name}" for p in players}
    player_id = st.selectbox("Ver ficha de", [p.id for p in players], format_func=lambda i: labels[i])
    player = next(p for p in players if p.id == player_id)
    render_player_card(session, team, player, last_n, season, league)


def _pdf_safe(text: object) -> str:
    """Sanea texto para el PDF.

    Las fuentes base de fpdf2 (Helvetica, aquí usada) solo soportan Latin-1;
    un nombre con caracteres fuera de ese rango (p.ej. "Žalgiris", "Šamanić")
    hace que `pdf.cell()` lance `FPDFUnicodeEncodingException` y reviente
    toda la generación del informe. Se aproximan esos caracteres a su
    equivalente ASCII más cercano (quitando diacríticos) y, si aun así queda
    algo fuera de Latin-1, se sustituye por "?" en vez de fallar.
    """
    without_accents = "".join(
        ch for ch in unicodedata.normalize("NFKD", str(text)) if not unicodedata.combining(ch)
    )
    return without_accents.encode("latin-1", "replace").decode("latin-1")


def _pdf_table(pdf: FPDF, df: pd.DataFrame) -> None:
    """Dibuja un DataFrame como tabla simple en el PDF."""
    if df.empty:
        pdf.set_font("Helvetica", "I", 9)
        pdf.cell(0, 6, "Sin datos.", new_x="LMARGIN", new_y="NEXT")
        return

    headers = list(df.columns)
    col_width = (pdf.w - pdf.l_margin - pdf.r_margin) / len(headers)
    pdf.set_font("Helvetica", "B", 8)
    for header in headers:
        pdf.cell(col_width, 6, _pdf_safe(header), border=1)
    pdf.ln()
    pdf.set_font("Helvetica", "", 8)
    for _, row in df.iterrows():
        for value in row:
            pdf.cell(col_width, 6, _pdf_safe(value), border=1)
        pdf.ln()


def build_pdf_report(
    session,
    team_a: models.Team,
    team_b: models.Team,
    last_n: int,
    season: "int | None" = None,
    league: "str | None" = None,
) -> bytes:
    """Genera el informe del enfrentamiento en PDF (mismos datos que la GUI)."""
    pdf = FPDF(orientation="L", unit="mm", format="A4")
    pdf.set_auto_page_break(auto=True, margin=12)
    pdf.add_page()

    pdf.set_font("Helvetica", "B", 16)
    pdf.cell(0, 10, _pdf_safe(f"Informe: {team_a.name} vs {team_b.name}"), new_x="LMARGIN", new_y="NEXT")
    pdf.set_font("Helvetica", "", 9)
    pdf.cell(0, 6, f"Generado el {datetime.now().strftime('%d/%m/%Y %H:%M')}", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(4)

    for team in (team_a, team_b):
        summary = team_advanced_summary(session, team, season=season, league=league)
        pdf.set_font("Helvetica", "B", 13)
        pdf.cell(0, 8, _pdf_safe(team.name), new_x="LMARGIN", new_y="NEXT")
        pdf.set_font("Helvetica", "", 9)
        pdf.cell(
            0,
            6,
            f"Pace: {_fmt(summary['avg_pace'])}  ORtg: {_fmt(summary['avg_off_rating'])}  "
            f"DRtg: {_fmt(summary['avg_def_rating'])}  Net: {_fmt(summary['avg_net_rating'])}  "
            f"eFG%: {_fmt_pct(summary['avg_efg_pct'])}  TS%: {_fmt_pct(summary['avg_ts_pct'])}",
            new_x="LMARGIN",
            new_y="NEXT",
        )
        pdf.ln(2)

        pdf.set_font("Helvetica", "B", 10)
        pdf.cell(0, 6, "Partidos guardados", new_x="LMARGIN", new_y="NEXT")
        _pdf_table(pdf, team_summary_df(session, team, season, league))
        pdf.ln(3)

        pdf.set_font("Helvetica", "B", 10)
        pdf.cell(0, 6, f"Forma reciente (ultimos {last_n} partidos)", new_x="LMARGIN", new_y="NEXT")
        _pdf_table(pdf, recent_form_df(session, team, last_n, season, league))
        pdf.ln(5)

    pdf.set_font("Helvetica", "B", 13)
    pdf.cell(0, 8, "Enfrentamientos directos", new_x="LMARGIN", new_y="NEXT")
    games = head_to_head_games(session, team_a, team_b, season, league)
    if not games:
        pdf.set_font("Helvetica", "I", 9)
        pdf.cell(0, 6, "Sin enfrentamientos directos guardados.", new_x="LMARGIN", new_y="NEXT")
    for game in games:
        pdf.set_font("Helvetica", "B", 11)
        pdf.cell(
            0,
            7,
            _pdf_safe(
                f"{format_date_es(game.date)} - {game.home_team.name} {game.home_score} - "
                f"{game.away_score} {game.away_team.name}"
            ),
            new_x="LMARGIN",
            new_y="NEXT",
        )
        for side_team in (game.home_team, game.away_team):
            pdf.set_font("Helvetica", "B", 9)
            pdf.cell(0, 6, _pdf_safe(side_team.name), new_x="LMARGIN", new_y="NEXT")
            _pdf_table(pdf, boxscore_df(session, game, side_team))
            pdf.ln(2)
        pdf.ln(3)

    pdf.set_font("Helvetica", "B", 13)
    pdf.cell(0, 8, "Avisos de calidad de datos", new_x="LMARGIN", new_y="NEXT")
    pdf.set_font("Helvetica", "", 9)
    warnings = validate_data(session)
    if not warnings:
        pdf.cell(0, 6, "Sin incidencias detectadas.", new_x="LMARGIN", new_y="NEXT")
    else:
        for warning in warnings:
            pdf.multi_cell(0, 6, _pdf_safe(f"- {warning}"))

    return bytes(pdf.output())


def main() -> None:
    session = get_session()
    focus_slug = config.TEAMS[0] if config.TEAMS else None
    baskonia = session.query(models.Team).filter_by(slug=focus_slug).first() if focus_slug else None
    if baskonia is None:
        st.error(
            f"No se encontró el equipo '{focus_slug}' en la base de datos. "
            "Ejecuta primero `python main.py`."
        )
        return

    seasons = list_seasons(session, baskonia)
    if not seasons:
        st.info("Sin temporadas con partidos registrados. Ejecuta `python main.py`.")
        return

    header_logo_col, header_title_col, header_n_col, header_season_col, header_league_col = st.columns(
        [1, 5, 2, 2, 2]
    )
    with header_logo_col:
        show_team_logo(baskonia.slug, width=64)
    with header_title_col:
        st.title(baskonia.name)
    with header_n_col:
        last_n = st.number_input("Últimos N partidos (forma)", min_value=1, max_value=20, value=5)
    with header_season_col:
        default_season = current_season(session, baskonia)
        season = st.selectbox(
            "Temporada",
            options=seasons,
            format_func=season_label,
            index=seasons.index(default_season) if default_season in seasons else 0,
            key="season_selector",
        )
    with header_league_col:
        # "Todas" (None) por defecto: ninguna competición es más relevante a priori,
        # y preseleccionar una ocultaría de entrada los partidos de las otras.
        league = st.selectbox(
            "Competición",
            options=[None] + list_leagues(session, baskonia),
            format_func=league_label,
            index=0,
            key="league_selector",
        )

    tab_resumen, tab_anteriores, tab_proximos, tab_plantilla = st.tabs(
        ["Resumen", "Partidos anteriores", "Próximos enfrentamientos", "Plantilla"]
    )
    with tab_resumen:
        render_team_tab(session, baskonia, last_n, season, league)
    with tab_anteriores:
        render_past_games_tab(session, baskonia, last_n, season, league)
    with tab_proximos:
        render_upcoming_tab(session, baskonia, last_n, season, league)
    with tab_plantilla:
        render_roster_tab(session, baskonia, last_n, season, league)


if __name__ == "__main__":
    main()
